"""
generate_slides.py
==================
Generates a 15-slide PowerPoint presentation for:

    "What Drives the Persistence in Current Account Imbalances?"
    Clower, Ito, Saadaoui — March 2026

Run from any directory:
    python code/slides/generate_slides.py

Output: slides/cab_persistency_slides.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT       = Path(__file__).resolve().parent.parent.parent   # updated_2025/
SLIDES_DIR = ROOT / "slides"
SLIDES_DIR.mkdir(exist_ok=True)
OUTPUT     = SLIDES_DIR / "cab_persistency_slides.pptx"

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
NAVY      = RGBColor(0x1F, 0x38, 0x64)   # dark navy — header bars, title bg
BLUE_ACC  = RGBColor(0x2E, 0x5B, 0xA8)   # medium blue — section dividers
GOLD      = RGBColor(0xC9, 0xA0, 0x27)   # gold accent rule
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF2, 0xF4, 0xF8)   # slide background
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x2E)   # near-black body text
GRAY_MID  = RGBColor(0x80, 0x80, 0x80)   # footnotes, slide numbers
HIGHLIGHT = RGBColor(0xE8, 0x4E, 0x1B)   # warm red for emphasis

# Slide geometry (16:9 widescreen)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)

# Header strip dimensions
HDR_H  = Inches(1.20)
MARGIN = Inches(0.45)
BODY_T = HDR_H + Inches(0.20)   # top of body area
BODY_H = SLIDE_H - BODY_T - Inches(0.55)
BODY_W = SLIDE_W - 2 * MARGIN

# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def _set_bg(slide, color: RGBColor):
    """Fill slide background with a solid color."""
    bg   = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, left, top, width, height,
          fill: RGBColor = None, line_width_pt: float = 0):
    """Add a rectangle auto-shape."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    shape.line.width = Pt(line_width_pt)
    if line_width_pt == 0:
        shape.line.fill.background()
    return shape


def _tb(slide, left, top, width, height):
    """Add a text box and return its text frame."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf  = box.text_frame
    tf.word_wrap = True
    return tf


def _para(tf, text: str, size_pt: int, bold: bool = False,
          color: RGBColor = TEXT_DARK, align=PP_ALIGN.LEFT,
          italic: bool = False, space_before: int = 0):
    """Add a paragraph to a text frame. Returns the paragraph."""
    p    = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run  = p.add_run()
    run.text = text
    run.font.size  = Pt(size_pt)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = "Calibri"
    return p


def _header_bar(slide, title: str, bg: RGBColor = NAVY,
                title_color: RGBColor = WHITE, size_pt: int = 24):
    """Add the navy header strip with slide title."""
    _rect(slide, Inches(0), Inches(0), SLIDE_W, HDR_H, fill=bg)
    tf = _tb(slide, MARGIN, Inches(0.22), SLIDE_W - 2 * MARGIN, HDR_H)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size  = Pt(size_pt)
    run.font.bold  = True
    run.font.color.rgb = title_color
    run.font.name  = "Calibri"


def _slide_number(slide, n: int, total: int = 15):
    tf = _tb(slide, SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.38),
             Inches(1.4), Inches(0.36))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{n} / {total}"
    run.font.size  = Pt(11)
    run.font.color.rgb = GRAY_MID
    run.font.name  = "Calibri"


def _bullet_list(slide, items: list, left=None, top=None,
                 width=None, height=None, base_size: int = 19):
    """
    items: list of (indent_level, text) or just strings for level-0.
    indent_level 0 = main bullet, 1 = sub-bullet.
    """
    if left   is None: left   = MARGIN
    if top    is None: top    = BODY_T
    if width  is None: width  = BODY_W
    if height is None: height = BODY_H

    tf = _tb(slide, left, top, width, height)
    first = True
    for item in items:
        if isinstance(item, str):
            level, text = 0, item
        else:
            level, text = item

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        indent = Inches(0.25 * level)
        p.space_before = Pt(4 if level == 0 else 2)

        # Bullet character
        bullet_char = "\u2022" if level == 0 else "\u2013"
        full_text   = f"  {bullet_char}  {text}"

        sz   = base_size - 2 * level
        col  = TEXT_DARK if level == 0 else GRAY_MID

        run = p.add_run()
        run.text = full_text
        run.font.size  = Pt(sz)
        run.font.color.rgb = col
        run.font.name  = "Calibri"
        p.alignment = PP_ALIGN.LEFT


def _content_slide(prs, slide_n: int, title: str, bg: RGBColor = LIGHT_BG):
    """Add a blank content slide with navy header and return (slide)."""
    layout = prs.slide_layouts[6]   # blank
    slide  = prs.slides.add_slide(layout)
    _set_bg(slide, bg)
    _header_bar(slide, title)
    _slide_number(slide, slide_n)
    return slide


def _section_divider(prs, section_label: str, slide_n: int):
    """Dark-navy full-slide section divider."""
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    _set_bg(slide, NAVY)

    # Gold rule at 40% height
    _rect(slide, Inches(1.5), Inches(3.1), Inches(10.33), Inches(0.06),
          fill=GOLD)

    tf = _tb(slide, Inches(1.5), Inches(2.2), Inches(10.33), Inches(0.85))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = section_label
    run.font.size  = Pt(32)
    run.font.bold  = True
    run.font.color.rgb = WHITE
    run.font.name  = "Calibri"

    _slide_number(slide, slide_n)
    return slide


# ---------------------------------------------------------------------------
# Individual slide builders
# ---------------------------------------------------------------------------

def slide_01_title(prs):
    """Title slide — dark navy, white text, gold rule."""
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    _set_bg(slide, NAVY)

    # Gold horizontal rule
    _rect(slide, Inches(0.8), Inches(2.95), Inches(11.73), Inches(0.07),
          fill=GOLD)

    # Main title
    tf = _tb(slide, Inches(0.8), Inches(1.2), Inches(11.73), Inches(1.7))
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "What Drives the Persistence\nin Current Account Imbalances?"
    run.font.size  = Pt(38)
    run.font.bold  = True
    run.font.color.rgb = WHITE
    run.font.name  = "Calibri"

    # Authors
    tf2 = _tb(slide, Inches(0.8), Inches(3.20), Inches(11.73), Inches(0.55))
    p2  = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = "Eric Clower  \u00b7  Hiro Ito  \u00b7  Jamel Saadaoui"
    run2.font.size  = Pt(22)
    run2.font.bold  = False
    run2.font.color.rgb = RGBColor(0xC8, 0xD8, 0xF0)
    run2.font.name  = "Calibri"

    # Affiliations
    tf3 = _tb(slide, Inches(0.8), Inches(3.85), Inches(11.73), Inches(0.55))
    p3  = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.LEFT
    run3 = p3.add_run()
    run3.text = ("Aptech Systems, Inc.  \u00b7  Portland State University  "
                 "\u00b7  Universit\u00e9 Paris 8")
    run3.font.size  = Pt(16)
    run3.font.italic = True
    run3.font.color.rgb = RGBColor(0xA0, 0xB8, 0xD8)
    run3.font.name  = "Calibri"

    # Conference / date line
    tf4 = _tb(slide, Inches(0.8), Inches(6.6), Inches(11.73), Inches(0.55))
    p4  = tf4.paragraphs[0]
    p4.alignment = PP_ALIGN.LEFT
    run4 = p4.add_run()
    run4.text = "March 2026"
    run4.font.size  = Pt(14)
    run4.font.color.rgb = GOLD
    run4.font.name  = "Calibri"

    return slide


def slide_02_motivation(prs):
    slide = _content_slide(prs, 2, "Motivation: Why Does Persistence Matter?")
    _bullet_list(slide, [
        "Global current account imbalances have grown in magnitude and dispersion since the 1980s",
        (1, "The cross-sectional variance of CA balances has roughly doubled since financial globalization accelerated"),
        "Persistent deficits and surpluses trigger serious policy concerns",
        (1, "Credit downgrades, speculative attacks on government bonds, fiscal policy constraints"),
        (1, "Even advanced economies affected \u2014 Moody\u2019s downgraded the US in May 2025"),
        "Standard theory is clear: the CA should revert to zero in the long run",
        (1, "Intertemporal budget constraint (LRBC) requires stationarity of the current account"),
        "Empirical tests produce contradictory and inconclusive results",
        (1, "Same country, different unit root tests, different conclusions"),
        "\u2192  What structural factors push a country toward explosive, persistent, or self-correcting dynamics?",
    ])
    return slide


def slide_03_approach(prs):
    slide = _content_slide(prs, 3, "Our Approach: A Two-Stage Empirical Strategy")

    # Stage 1 box
    box1 = _rect(slide, MARGIN, BODY_T, Inches(5.8), Inches(2.5),
                 fill=RGBColor(0xE8, 0xED, 0xF7))
    tf1  = _tb(slide, MARGIN + Inches(0.15), BODY_T + Inches(0.12),
               Inches(5.5), Inches(2.3))
    _para(tf1, "Stage 1 \u2014 Regime Classification", 16, bold=True, color=NAVY)
    _para(tf1, "Estimate a Markov-switching AR model for each country", 15, color=TEXT_DARK)
    _para(tf1, "\u2192  Classify each country-year observation as:", 15, color=TEXT_DARK)
    _para(tf1, "      Explosive  |  Unit root  |  Stationary", 15, bold=True, color=BLUE_ACC)
    _para(tf1, "74 countries, annual data, 1985\u20132023", 14, color=GRAY_MID, italic=True)

    # Arrow
    tf_arr = _tb(slide, Inches(6.5), BODY_T + Inches(0.85), Inches(0.9), Inches(0.9))
    p = tf_arr.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "\u2794"
    run.font.size = Pt(36)
    run.font.color.rgb = GOLD
    run.font.name = "Calibri"

    # Stage 2 box
    box2 = _rect(slide, Inches(7.5), BODY_T, Inches(5.35), Inches(2.5),
                 fill=RGBColor(0xE8, 0xED, 0xF7))
    tf2  = _tb(slide, Inches(7.65), BODY_T + Inches(0.12),
               Inches(5.0), Inches(2.3))
    _para(tf2, "Stage 2 \u2014 Structural Determinants", 16, bold=True, color=NAVY)
    _para(tf2, "Estimate a multinomial logit (MNL) model", 15, color=TEXT_DARK)
    _para(tf2, "\u2192  Outcome: regime from Stage 1", 15, color=TEXT_DARK)
    _para(tf2, "\u2192  Explain using structural country characteristics", 15, color=TEXT_DARK)
    _para(tf2, "Cluster-robust SEs, country-level clusters", 14, color=GRAY_MID, italic=True)

    # Bottom note
    tf_note = _tb(slide, MARGIN, BODY_T + Inches(3.0), BODY_W, Inches(0.5))
    _para(tf_note,
          "Key contribution: two-stage design separates time-series characterization from cross-sectional "
          "policy analysis \u2014 regime classifications feed directly into the multinomial model.",
          14, color=GRAY_MID, italic=True)
    return slide


def slide_04_facts(prs):
    slide = _content_slide(prs, 4,
        "Fact: Global CA Imbalances Are Larger and More Dispersed")
    _bullet_list(slide, [
        "The mean absolute current account balance (% GDP) has roughly doubled since the early 1980s",
        (1, "Financial globalization has expanded access to external financing, relaxing short-run adjustment"),
        "The distribution of CA balances has widened over time",
        (1, "Median is relatively stable; the tails have grown fatter \u2014 more extreme surpluses and deficits"),
        "Sharp cyclical spikes around the GFC (2008\u201309) and COVID (2020\u201321) recessions",
        (1, "Followed by partial reversion, but not full convergence"),
        "Cross-country heterogeneity is pronounced",
        (1, "Advanced: moderate, stable imbalances"),
        (1, "Emerging: larger cyclical swings, higher explosive frequencies"),
        (1, "Less-developed: most volatile, least mean-reverting"),
        "\u2192  Persistence is the norm; the extent varies systematically with country characteristics",
    ])
    # figure reference
    tf = _tb(slide, MARGIN, SLIDE_H - Inches(0.5), BODY_W, Inches(0.38))
    _para(tf, "Figures 1\u20133 in paper", 12, color=GRAY_MID, italic=True)
    return slide


def slide_05_theory(prs):
    slide = _content_slide(prs, 5, "Theory: Sustainability Requires Stationarity")
    _bullet_list(slide, [
        "The intertemporal approach provides the standard benchmark (Trehan & Walsh 1991; Taylor 2002)",
        (1, "LRBC: current account stationarity (as % of GDP) is a necessary condition for external sustainability"),
        (1, "Unit root \u2192 permanent shocks; Explosive \u2192 diverging dynamics"),
        "But: stationarity is a long-run property \u2014 episodes of local non-stationarity may be consistent "
        "with global sustainability",
        (1, "A CA may be locally explosive during a crisis yet globally sustainable if episodes are transient"),
        (1, "Markov-switching framework handles this directly via global stationarity conditions "
            "(Psaradakis et al. 2004)"),
        "Bohn (1998) fiscal reaction-function approach: sustainability holds if the primary surplus responds "
        "positively to debt accumulation",
        (1, "Complements the unit root perspective: both require that adjustment mechanisms are active"),
        "\u2192  We need a model that allows dynamics to differ across time periods within each country",
    ])
    return slide


def slide_06_linear_tests(prs):
    slide = _content_slide(prs, 6, "Problem: Linear Unit Root Tests Give Inconsistent Results")

    _bullet_list(slide, [
        "Standard linear tests (ADF, DFGLS, Phillips-Perron) applied to CA/GDP across 74 countries",
        "Results are contradictory and sensitive to the choice of test",
        (1, "ADF and DFGLS: roughly half of countries cannot reject the unit root null"),
        (1, "Phillips-Perron: many more rejections \u2014 same country, opposite conclusion"),
        "Root cause: linear tests assume constant AR parameters across time",
        (1, "When persistence or volatility shifts (policy changes, crises), "
            "linear tests confound regime changes with structural non-stationarity"),
        "We treat these results as a descriptive baseline, not a verdict",
        "\u2192  The correct framework must allow dynamics to change over time",
    ])
    # figure reference
    tf = _tb(slide, MARGIN, SLIDE_H - Inches(0.5), BODY_W, Inches(0.38))
    _para(tf, "Figure 4 in paper | Appendix 2: country-by-country results", 12,
          color=GRAY_MID, italic=True)
    return slide


def slide_07_ms_model(prs):
    slide = _content_slide(prs, 7, "Stage 1: Markov-Switching AR Model")
    _bullet_list(slide, [
        "MS-AR(1) model with independently switching intercept (\u03bc), "
        "AR coefficient (\u03c1), and variance (\u03c3\u00b2)",
        (1, "Each component follows an independent two-state Markov chain"),
        (1, "Allows persistence and volatility to evolve separately \u2014 critical for robust inference"),
        "Three regime types, classified by estimated persistence parameter \u03c1:",
    ])

    # Three regime boxes
    y_boxes = BODY_T + Inches(2.1)
    w_box = Inches(3.7)
    specs = [
        (MARGIN,              "Stationary",  "|\\u03c1| < 1",  RGBColor(0x27, 0x7A, 0x3C)),
        (MARGIN + Inches(4.0), "Unit Root",   "\\u03c1 \u2248 1", RGBColor(0x1F, 0x38, 0x64)),
        (MARGIN + Inches(8.0), "Explosive",   "|\\u03c1| > 1",  HIGHLIGHT),
    ]
    for xpos, label, rho_text, col in specs:
        _rect(slide, xpos, y_boxes, w_box, Inches(1.35), fill=col)
        tf = _tb(slide, xpos + Inches(0.12), y_boxes + Inches(0.12),
                 w_box - Inches(0.24), Inches(1.1))
        _para(tf, label, 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _para(tf, rho_text.replace("\\u03c1", "\u03c1"), 15, color=WHITE,
              align=PP_ALIGN.CENTER)

    # Footer bullets
    tf_f = _tb(slide, MARGIN, y_boxes + Inches(1.55), BODY_W, Inches(1.5))
    _para(tf_f, "  \u2022  Parametric bootstrap hypothesis tests (1,000 replications) for each country",
          15, color=TEXT_DARK)
    _para(tf_f, "  \u2022  Global stationarity assessed via second-order conditions (Psaradakis et al. 2004)",
          15, color=TEXT_DARK)
    _para(tf_f, "  \u2022  Reference: Morita, Psaradakis, Sola & Yunis (2024); Hall, Psaradakis & Sola (1999)",
          14, color=GRAY_MID, italic=True)
    return slide


def slide_08_regime_dist(prs):
    slide = _content_slide(prs, 8,
        "Stage 1 Results: Most Country-Years Are in the Unit Root Regime")

    # Three big stat callouts
    y_stat = BODY_T + Inches(0.1)
    stats  = [
        (MARGIN,               "57%",  "Unit Root",   NAVY),
        (MARGIN + Inches(4.3), "38%",  "Stationary",  RGBColor(0x27, 0x7A, 0x3C)),
        (MARGIN + Inches(8.6), "5%",   "Explosive",   HIGHLIGHT),
    ]
    for xpos, pct, label, col in stats:
        _rect(slide, xpos, y_stat, Inches(3.7), Inches(1.7), fill=col)
        tf = _tb(slide, xpos + Inches(0.1), y_stat + Inches(0.08),
                 Inches(3.5), Inches(1.55))
        _para(tf, pct,   38, bold=True,  color=WHITE, align=PP_ALIGN.CENTER)
        _para(tf, label, 18, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    _bullet_list(slide, [
        "Unit root regime is the modal state: CA dynamics are highly persistent but not diverging",
        "Explosive episodes are episodic, not chronic \u2014 consistent with temporary adjustment crises",
        (1, "Concentrated around the GFC (2008\u201309), COVID (2020\u201321), and EM stress periods"),
        "Stationary behavior is substantial (~38%) \u2014 many country-years do exhibit mean reversion",
        "\u2192  Significant cross-country and time variation to explain in Stage 2",
    ], top=y_stat + Inches(1.95))
    tf = _tb(slide, MARGIN, SLIDE_H - Inches(0.5), BODY_W, Inches(0.38))
    _para(tf, "Figure 5 in paper | Quarterly observations", 12, color=GRAY_MID, italic=True)
    return slide


def slide_09_heterogeneity(prs):
    slide = _content_slide(prs, 9,
        "Regime Dynamics Differ Systematically Across Country Groups")

    # Two-column layout: income groups | time evolution
    col_w = Inches(5.8)
    col2_x = MARGIN + col_w + Inches(0.4)

    # Left column header
    tf_lh = _tb(slide, MARGIN, BODY_T, col_w, Inches(0.4))
    _para(tf_lh, "Income Groups", 16, bold=True, color=NAVY)

    rows = [
        ("Advanced economies",  "54% UnitR  |  45% Stat  |   1% Expl"),
        ("Emerging markets",    "56% UnitR  |  32% Stat  |  12% Expl"),
        ("Less-developed",      "Similar to emerging; elevated explosive share"),
    ]
    y_row = BODY_T + Inches(0.50)
    row_h = Inches(0.85)
    for i, (grp, detail) in enumerate(rows):
        bg = RGBColor(0xE8, 0xED, 0xF7) if i % 2 == 0 else WHITE
        _rect(slide, MARGIN, y_row + i * row_h, col_w, row_h - Inches(0.05), fill=bg)
        tf = _tb(slide, MARGIN + Inches(0.1), y_row + i * row_h + Inches(0.05),
                 col_w - Inches(0.2), row_h - Inches(0.1))
        _para(tf, grp,    14, bold=True,  color=NAVY)
        _para(tf, detail, 13, bold=False, color=TEXT_DARK)

    # Thin separator
    _rect(slide, col2_x - Inches(0.2), BODY_T, Inches(0.03),
          BODY_H, fill=RGBColor(0xCC, 0xCC, 0xCC))

    # Right column header
    tf_rh = _tb(slide, col2_x, BODY_T, col_w, Inches(0.4))
    _para(tf_rh, "Time Evolution", 16, bold=True, color=NAVY)

    tf_r = _tb(slide, col2_x, BODY_T + Inches(0.50), col_w, BODY_H - Inches(0.5))
    _para(tf_r, "  \u2022  Explosive episodes cluster around identifiable crises", 15, color=TEXT_DARK)
    _para(tf_r, "       GFC (2008\u201309)  |  Euro crisis (2010\u201312)  |  COVID (2020\u201321)", 14,
          color=GRAY_MID)
    _para(tf_r, "  \u2022  Unit root share rises steadily post-GFC", 15, color=TEXT_DARK)
    _para(tf_r, "  \u2022  Regional patterns:", 15, color=TEXT_DARK)
    _para(tf_r, "       Latin America: highest unit root share (62%)", 14, color=GRAY_MID)
    _para(tf_r, "       East Asia: comparatively stable, lower explosive frequency", 14, color=GRAY_MID)
    _para(tf_r, "  \u2022  Advanced economies almost never enter explosive regime", 15, color=TEXT_DARK)

    tf = _tb(slide, MARGIN, SLIDE_H - Inches(0.5), BODY_W, Inches(0.38))
    _para(tf, "Figures 6\u20137 in paper", 12, color=GRAY_MID, italic=True)
    return slide


def slide_10_mnl_setup(prs):
    slide = _content_slide(prs, 10, "Stage 2: Multinomial Logit (MNL) Model")

    # Outcome box
    _rect(slide, MARGIN, BODY_T, BODY_W, Inches(0.72), fill=RGBColor(0xE0, 0xE8, 0xF5))
    tf_out = _tb(slide, MARGIN + Inches(0.15), BODY_T + Inches(0.08),
                 BODY_W - Inches(0.3), Inches(0.6))
    _para(tf_out,
          "Outcome: regime classification from Stage 1  "
          "\u2192  Explosive (reference)  |  Stationary  |  Unit root",
          16, bold=True, color=NAVY)

    # Six covariates
    tf_cov = _tb(slide, MARGIN, BODY_T + Inches(0.9), BODY_W, Inches(0.35))
    _para(tf_cov, "Six structural regressors:", 15, bold=True, color=NAVY)

    cov_rows = [
        ("Fixed exchange rate (dummy)",
         "1 if fixed/pegged arrangement in year t  (Ilzetzki, Reinhart & Rogoff)"),
        ("Current account deficit (dummy)",
         "1 if country ran a CA deficit in year t"),
        ("Commodity exporter (dummy)",
         "1 if primary commodity exporter (terms-of-trade volatility exposure)"),
        ("Manufacturing exporter (dummy)",
         "1 if manufactured goods are the primary export category"),
        ("Trade openness (lagged, t\u22121)",
         "Total trade / GDP, lagged one period to reduce simultaneity"),
        ("Fiscal balance (3-yr avg, lagged)",
         "General gov\u2019t balance / GDP, 3-yr moving avg, lagged one period"),
    ]
    y0 = BODY_T + Inches(1.3)
    rh = Inches(0.52)
    for i, (var, desc) in enumerate(cov_rows):
        bg = RGBColor(0xF5, 0xF7, 0xFA) if i % 2 == 0 else WHITE
        _rect(slide, MARGIN, y0 + i * rh, BODY_W, rh - Inches(0.02), fill=bg)
        tf = _tb(slide, MARGIN + Inches(0.1), y0 + i * rh + Inches(0.04),
                 BODY_W - Inches(0.2), rh - Inches(0.06))
        _para(tf, f"{var}:  {desc}", 14, color=TEXT_DARK)

    tf_f = _tb(slide, MARGIN, y0 + 6 * rh + Inches(0.05), BODY_W, Inches(0.38))
    _para(tf_f,
          "74 countries \u00b7 2,451 observations \u00b7 cluster-robust SEs (country level) "
          "\u00b7 reference category: Explosive",
          13, color=GRAY_MID, italic=True)
    return slide


def slide_11_baseline_results(prs):
    slide = _content_slide(prs, 11,
        "Baseline Results: Fixed ER and Manufacturing Exports Are Key Predictors")

    # Simplified Table 4 summary
    tf_intro = _tb(slide, MARGIN, BODY_T, BODY_W, Inches(0.4))
    _para(tf_intro,
          "MNL coefficients (explosive = reference) \u2014 Table 4 in paper",
          15, bold=True, color=NAVY)

    # Header row
    hdr_y = BODY_T + Inches(0.45)
    _rect(slide, MARGIN, hdr_y, BODY_W, Inches(0.38), fill=NAVY)
    hdrs = ["Regressor", "Stationary vs Explosive", "Unit Root vs Explosive"]
    col_xs = [MARGIN + Inches(0.1), MARGIN + Inches(4.7), MARGIN + Inches(8.4)]
    tf_h = _tb(slide, MARGIN + Inches(0.1), hdr_y + Inches(0.04),
               BODY_W - Inches(0.2), Inches(0.32))
    p = tf_h.paragraphs[0]
    run = p.add_run()
    run.text = (f"{'Regressor':<36}  {'Stationary vs Explosive':^22}  "
                f"{'Unit Root vs Explosive':^22}")
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"

    rows = [
        ("Fixed exchange rate",           "+  *** (large)",  "+  *** (large)"),
        ("Manufacturing exporter",        "+  ***",          "+  **"),
        ("Trade openness (t\u22121)",     "\u2212  **",       "\u2212  *"),
        ("Current account deficit",       "n.s.",             "n.s."),
        ("Commodity exporter",            "n.s.",             "n.s."),
        ("Fiscal balance (t\u22121)",     "n.s.",             "n.s."),
    ]
    row_h = Inches(0.52)
    for i, (var, stat, ur) in enumerate(rows):
        bg = RGBColor(0xF2, 0xF4, 0xF8) if i % 2 == 0 else WHITE
        y_r = hdr_y + Inches(0.38) + i * row_h
        _rect(slide, MARGIN, y_r, BODY_W, row_h, fill=bg)
        tf_r = _tb(slide, MARGIN + Inches(0.1), y_r + Inches(0.09),
                   BODY_W - Inches(0.2), row_h - Inches(0.1))
        p_r = tf_r.paragraphs[0]
        run_r = p_r.add_run()
        run_r.text = f"{var:<34}  {stat:^24}  {ur:^24}"
        run_r.font.size = Pt(14)
        run_r.font.color.rgb = TEXT_DARK
        run_r.font.name = "Calibri"

    tf_f = _tb(slide, MARGIN, SLIDE_H - Inches(0.5), BODY_W, Inches(0.38))
    _para(tf_f,
          "*p<0.10  **p<0.05  ***p<0.01  \u00b7  cluster-robust SEs  \u00b7  "
          "n.s. = not significant",
          12, color=GRAY_MID, italic=True)
    return slide


def slide_12_ames(prs):
    slide = _content_slide(prs, 12,
        "Average Marginal Effects: Fixed ER Reduces Explosive Probability by ~9 pp")

    # Three callout boxes for AMEs
    y_boxes = BODY_T + Inches(0.1)
    ame_specs = [
        (MARGIN,               "\u22129.1 pp",
         "Fixed ER \u2192 Explosive prob.", NAVY),
        (MARGIN + Inches(4.5), "+5.9 pp",
         "Fixed ER \u2192 Stationary prob.", RGBColor(0x27, 0x7A, 0x3C)),
        (MARGIN + Inches(9.0), "+3.2 pp",
         "Fixed ER \u2192 Unit root prob.",  BLUE_ACC),
    ]
    for xpos, val, label, col in ame_specs:
        _rect(slide, xpos, y_boxes, Inches(3.85), Inches(1.45), fill=col)
        tf = _tb(slide, xpos + Inches(0.12), y_boxes + Inches(0.08),
                 Inches(3.6), Inches(1.3))
        _para(tf, val,   28, bold=True,  color=WHITE, align=PP_ALIGN.CENTER)
        _para(tf, label, 14, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    _bullet_list(slide, [
        "Fixed exchange rate is the dominant predictor \u2014 reduces explosive probability by ~9 percentage points",
        (1, "Consistent with the disciplining role of exchange rate pegs: "
            "defending the peg forces external adjustment"),
        "Manufacturing exporter: significant negative effect on explosive probability",
        "No other covariate generates a statistically significant marginal effect on any regime",
        "\u2192  KEY ASYMMETRY: Structural variables explain explosive vs. non-explosive propensity,",
        (1, "but do NOT distinguish stationary from unit root behavior within the stable regime space"),
        (1, "Binary logit (stationary vs. unit root only): all 6 regressors insignificant (p > 0.25)"),
    ], top=y_boxes + Inches(1.65))
    tf = _tb(slide, MARGIN, SLIDE_H - Inches(0.5), BODY_W, Inches(0.38))
    _para(tf, "Table 5 (AMEs) in paper \u00b7 statistical significance based on cluster-robust SEs", 12,
          color=GRAY_MID, italic=True)
    return slide


def slide_13_robustness(prs):
    slide = _content_slide(prs, 13, "Robustness Checks")
    _bullet_list(slide, [
        "Lagged exchange rate regime (endogeneity concern)",
        (1, "Re-estimate with fixed_er lagged one period \u2192 results unchanged"),
        (1, "Addresses concern that ER regime is simultaneously determined with CA dynamics"),
        "Year fixed effects and global crisis dummies",
        (1, "Models with annual time dummies: consistently higher BIC \u2014 global shocks do not dominate"),
        (1, "Explicit GFC (2008\u201309) and COVID (2020\u201321) dummies: do not improve fit"),
        "Oil exporters",
        (1, "Narrow oil exporter dummy covers only 40 observations \u2014 too few for full MNL"),
        (1, "Regime distribution is distinctive: 33% explosive, 68% stationary, 0% unit root"),
        "Cross-validation (country-grouped 5-fold)",
        (1, "Countries (not rows) assigned to folds \u2014 prevents data leakage across country time series"),
        (1, "Baseline model achieves solid out-of-sample log-loss; lagged ER performs similarly"),
        "Binary logit (stationary vs. unit root)",
        (1, "Drops explosive observations; all 6 regressors insignificant (p > 0.25) \u2014 confirms asymmetry"),
    ])
    return slide


def slide_14_subgroups(prs):
    slide = _content_slide(prs, 14,
        "Subgroup Results: Emerging Markets and Post-GFC Drive the Findings")

    # Two columns
    col_w = Inches(5.8)
    col2_x = MARGIN + col_w + Inches(0.4)

    # Left: income groups
    tf_lh = _tb(slide, MARGIN, BODY_T, col_w, Inches(0.4))
    _para(tf_lh, "By Income Group", 16, bold=True, color=NAVY)

    income_rows = [
        ("Full sample",      "0.050", "2,451 obs"),
        ("Advanced",         "n/a *",  "1,041 obs"),
        ("Emerging markets", "0.186",  "903 obs"),
        ("Low income",       "\u2014 (skipped)",   "< 5 explosive obs"),
    ]
    y0 = BODY_T + Inches(0.5)
    rh = Inches(0.75)
    for i, (grp, r2, note) in enumerate(income_rows):
        bg = RGBColor(0xE8, 0xED, 0xF7) if i % 2 == 0 else WHITE
        _rect(slide, MARGIN, y0 + i * rh, col_w, rh - Inches(0.04), fill=bg)
        col_r2 = RGBColor(0x27, 0x7A, 0x3C) if grp == "Emerging markets" else NAVY
        tf = _tb(slide, MARGIN + Inches(0.1), y0 + i * rh + Inches(0.06),
                 col_w - Inches(0.2), rh - Inches(0.08))
        _para(tf, grp,  14, bold=True,  color=NAVY)
        _para(tf, f"McFadden R\u00b2 = {r2}  \u00b7  {note}", 13, color=col_r2)

    tf_l_note = _tb(slide, MARGIN, y0 + 4 * rh + Inches(0.05), col_w, Inches(0.55))
    _para(tf_l_note,
          "* Advanced: only 15 explosive obs \u2014 R\u00b2 unreliable, marked n/a in paper.\n"
          "Emerging markets: ~4\u00d7 improvement over full sample.",
          13, color=GRAY_MID, italic=True)

    # Separator
    _rect(slide, col2_x - Inches(0.2), BODY_T, Inches(0.03),
          BODY_H, fill=RGBColor(0xCC, 0xCC, 0xCC))

    # Right: time periods
    tf_rh = _tb(slide, col2_x, BODY_T, col_w, Inches(0.4))
    _para(tf_rh, "By Time Period", 16, bold=True, color=NAVY)

    time_rows = [
        ("Full sample",  "0.050", "1985\u20132023"),
        ("Pre-GFC",      "0.042", "Before 2008"),
        ("Post-GFC",     "0.089", "2008 onward"),
    ]
    for i, (period, r2, span) in enumerate(time_rows):
        bg = RGBColor(0xE8, 0xED, 0xF7) if i % 2 == 0 else WHITE
        _rect(slide, col2_x, y0 + i * rh, col_w, rh - Inches(0.04), fill=bg)
        col_r2 = RGBColor(0x27, 0x7A, 0x3C) if period == "Post-GFC" else NAVY
        tf = _tb(slide, col2_x + Inches(0.1), y0 + i * rh + Inches(0.06),
                 col_w - Inches(0.2), rh - Inches(0.08))
        _para(tf, f"{period}  ({span})", 14, bold=True, color=NAVY)
        _para(tf, f"McFadden R\u00b2 = {r2}", 13, color=col_r2)

    tf_r_note = _tb(slide, col2_x, y0 + 3 * rh + Inches(0.05), col_w, Inches(0.9))
    _para(tf_r_note,
          "Post-GFC model fit more than doubles (R\u00b2: 0.042 \u2192 0.089).\n"
          "Structural policy variables have become increasingly relevant predictors "
          "of CA regime membership since the Global Financial Crisis.",
          13, color=GRAY_MID, italic=True)
    return slide


def slide_15_conclusion(prs):
    slide = _content_slide(prs, 15, "Conclusions")

    tf_lead = _tb(slide, MARGIN, BODY_T, BODY_W, Inches(0.4))
    _para(tf_lead, "Three principal findings:", 16, bold=True, color=NAVY)

    findings = [
        ("1",
         "Fixed exchange rate arrangements substantially reduce explosive CA dynamics",
         "Reducing explosive regime probability by ~9 percentage points on average. "
         "Fixed pegs act as a policy commitment device that constrains destabilizing dynamics."),
        ("2",
         "Structural variables explain explosive vs. non-explosive propensity \u2014 not within-stable-space behavior",
         "The same covariates that predict explosive entry do not distinguish stationary from unit root "
         "behavior. The model is an explosive propensity classifier."),
        ("3",
         "Results concentrate in emerging markets and have strengthened post-GFC",
         "McFadden R\u00b2 rises from 0.05 to 0.19 for EMs. Post-GFC R\u00b2 doubles. "
         "Exchange rate and trade structure are most consequential for middle-income countries."),
    ]
    y0 = BODY_T + Inches(0.55)
    fh = Inches(1.45)
    for i, (num, headline, detail) in enumerate(findings):
        _rect(slide, MARGIN, y0 + i * fh, Inches(0.55), fh - Inches(0.08), fill=NAVY)
        tf_num = _tb(slide, MARGIN + Inches(0.05), y0 + i * fh + Inches(0.30),
                     Inches(0.45), Inches(0.6))
        _para(tf_num, num, 24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        _rect(slide, MARGIN + Inches(0.57), y0 + i * fh,
              BODY_W - Inches(0.57), fh - Inches(0.08),
              fill=RGBColor(0xF2, 0xF4, 0xF8))
        tf_c = _tb(slide, MARGIN + Inches(0.72), y0 + i * fh + Inches(0.07),
                   BODY_W - Inches(0.82), fh - Inches(0.12))
        _para(tf_c, headline, 15, bold=True,  color=NAVY)
        _para(tf_c, detail,   13, bold=False, color=TEXT_DARK)

    tf_pol = _tb(slide, MARGIN, y0 + 3 * fh + Inches(0.05), BODY_W, Inches(0.55))
    _para(tf_pol,
          "Policy implication: ER regime choice remains central for external adjustment, "
          "particularly in emerging markets. Countries with fixed pegs face stronger incentives "
          "to adjust before imbalances become explosive.",
          14, color=GRAY_MID, italic=True)
    return slide


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def build_presentation():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building slides...")
    slide_01_title(prs)
    print("  1  Title slide")
    slide_02_motivation(prs)
    print("  2  Motivation")
    slide_03_approach(prs)
    print("  3  Two-stage approach")
    slide_04_facts(prs)
    print("  4  Facts")
    slide_05_theory(prs)
    print("  5  Theory")
    slide_06_linear_tests(prs)
    print("  6  Linear test limitations")
    slide_07_ms_model(prs)
    print("  7  Markov-switching model")
    slide_08_regime_dist(prs)
    print("  8  Regime distribution")
    slide_09_heterogeneity(prs)
    print("  9  Cross-country heterogeneity")
    slide_10_mnl_setup(prs)
    print("  10 MNL setup")
    slide_11_baseline_results(prs)
    print("  11 Baseline results")
    slide_12_ames(prs)
    print("  12 Average marginal effects")
    slide_13_robustness(prs)
    print("  13 Robustness checks")
    slide_14_subgroups(prs)
    print("  14 Subgroup heterogeneity")
    slide_15_conclusion(prs)
    print("  15 Conclusions")

    prs.save(OUTPUT)
    print(f"\nSaved: {OUTPUT}")
    print(f"Slides: 15  |  {prs.slide_width.inches:.2f}\u201d \u00d7 "
          f"{prs.slide_height.inches:.2f}\u201d (16:9)")


if __name__ == "__main__":
    build_presentation()
